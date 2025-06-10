from flask import Flask, request, send_file, jsonify
from PIL import Image
import io
import os
from flask_cors import CORS
from werkzeug.utils import secure_filename
import tempfile
import zipfile
from datetime import datetime

# PDF and document processing libraries
try:
    import PyPDF2
    from PyPDF2 import PdfWriter, PdfReader
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    import fitz  # PyMuPDF
    from docx import Document
    from docx2pdf import convert as docx_to_pdf_convert
    import pandas as pd
    from openpyxl import Workbook
    import pptx
    from pptx import Presentation
    import cv2
    import numpy as np
    from fpdf import FPDF
    import pdfkit
    import base64
    from cryptography.hazmat.primitives import hashes
    from cryptography.hazmat.primitives.asymmetric import rsa, padding
    from cryptography.hazmat.primitives.serialization import Encoding, PrivateFormat, NoEncryption
    import qrcode
except ImportError as e:
    print(f"Some libraries are missing. Install with: pip install {e.name}")

app = Flask(__name__)
CORS(app)

# Configuration
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
UPLOAD_FOLDER = tempfile.gettempdir()
ALLOWED_EXTENSIONS = {
    'pdf', 'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx', 
    'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'html', 'htm'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_temp_filename(extension):
    return os.path.join(UPLOAD_FOLDER, f"temp_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.{extension}")

# Original image conversion endpoint
@app.route('/convert-image', methods=['POST'])
def convert_image():
    if 'image' not in request.files:
        return jsonify({'error': 'No image uploaded'}), 400

    image_file = request.files['image']
    target_format = request.form.get('format')

    if not target_format:
        return jsonify({'error': 'No target format specified'}), 400

    try:
        img = Image.open(image_file)

        if target_format.lower() in ['jpeg', 'jpg'] and img.mode in ("RGBA", "P"):
            img = img.convert("RGB")

        img_bytes = io.BytesIO()
        img.save(img_bytes, format=target_format.upper())
        img_bytes.seek(0)

        mime_type = f'image/{target_format.lower()}'
        return send_file(img_bytes, mimetype=mime_type, download_name=f'converted.{target_format.lower()}')

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 1. Merge PDF
@app.route('/merge-pdf', methods=['POST'])
def merge_pdf():
    if 'files' not in request.files:
        return jsonify({'error': 'No PDF files uploaded'}), 400
    
    files = request.files.getlist('files')
    if len(files) < 2:
        return jsonify({'error': 'At least 2 PDF files required for merging'}), 400
    
    try:
        merger = PdfWriter()
        
        for file in files:
            if file and file.filename.endswith('.pdf'):
                reader = PdfReader(file)
                for page in reader.pages:
                    merger.add_page(page)
        
        output = io.BytesIO()
        merger.write(output)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='merged.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 2. Split PDF
@app.route('/split-pdf', methods=['POST'])
def split_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    split_type = request.form.get('split_type', 'all')  # 'all', 'range', 'single'
    page_range = request.form.get('page_range', '')
    
    try:
        reader = PdfReader(file)
        total_pages = len(reader.pages)
        
        if split_type == 'all':
            # Split into individual pages
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
                for i, page in enumerate(reader.pages):
                    writer = PdfWriter()
                    writer.add_page(page)
                    
                    page_buffer = io.BytesIO()
                    writer.write(page_buffer)
                    page_buffer.seek(0)
                    
                    zip_file.writestr(f'page_{i+1}.pdf', page_buffer.read())
            
            zip_buffer.seek(0)
            return send_file(zip_buffer, mimetype='application/zip', download_name='split_pages.zip')
        
        elif split_type == 'single':
            page_num = int(request.form.get('page_number', 1)) - 1
            if 0 <= page_num < total_pages:
                writer = PdfWriter()
                writer.add_page(reader.pages[page_num])
                
                output = io.BytesIO()
                writer.write(output)
                output.seek(0)
                
                return send_file(output, mimetype='application/pdf', download_name=f'page_{page_num+1}.pdf')
        
        return jsonify({'error': 'Invalid split parameters'}), 400
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 3. Compress PDF
@app.route('/compress-pdf', methods=['POST'])
def compress_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    compression_level = request.form.get('compression_level', 'medium')
    
    try:
        # Save uploaded file temporarily
        temp_input = get_temp_filename('pdf')
        file.save(temp_input)
        
        # Open with PyMuPDF for compression
        doc = fitz.open(temp_input)
        
        # Compression settings based on level
        if compression_level == 'high':
            deflate_level = 9
            image_quality = 50
        elif compression_level == 'low':
            deflate_level = 3
            image_quality = 85
        else:  # medium
            deflate_level = 6
            image_quality = 70
        
        # Compress images in PDF
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("jpeg", jpg_quality=image_quality)
                    doc._replace_image(xref, img_data)
                
                pix = None
        
        # Save compressed PDF
        temp_output = get_temp_filename('pdf')
        doc.save(temp_output, deflate=True, deflate_level=deflate_level)
        doc.close()
        
        # Clean up
        os.remove(temp_input)
        
        return send_file(temp_output, mimetype='application/pdf', download_name='compressed.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 4. PDF to Word
@app.route('/pdf-to-word', methods=['POST'])
def pdf_to_word():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        # Extract text from PDF
        reader = PdfReader(file)
        doc = Document()
        
        for page in reader.pages:
            text = page.extract_text()
            if text.strip():
                doc.add_paragraph(text)
                doc.add_page_break()
        
        # Save to BytesIO
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        
        return send_file(output, 
                        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                        download_name='converted.docx')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 5. PDF to PowerPoint
@app.route('/pdf-to-powerpoint', methods=['POST'])
def pdf_to_powerpoint():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        reader = PdfReader(file)
        prs = Presentation()
        
        for page in reader.pages:
            text = page.extract_text()
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = f"Page {len(prs.slides)}"
            content.text = text[:500] + "..." if len(text) > 500 else text
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return send_file(output, 
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        download_name='converted.pptx')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 6. PDF to Excel
@app.route('/pdf-to-excel', methods=['POST'])
def pdf_to_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        reader = PdfReader(file)
        wb = Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        
        row = 1
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            ws.cell(row=row, column=1, value=f"Page {page_num + 1}")
            row += 1
            
            # Split text into lines and add to cells
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    ws.cell(row=row, column=1, value=line.strip())
                    row += 1
            row += 1  # Add spacing between pages
        
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)
        
        return send_file(output, 
                        mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                        download_name='converted.xlsx')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 7. Word to PDF
@app.route('/word-to-pdf', methods=['POST'])
def word_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No Word file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        # Read Word document
        doc = Document(file)
        
        # Create PDF using FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Handle text encoding
                text = paragraph.text.encode('latin1', 'replace').decode('latin1')
                pdf.multi_cell(0, 10, text)
                pdf.ln(2)
        
        output = io.BytesIO()
        pdf_content = pdf.output(dest='S').encode('latin1')
        output.write(pdf_content)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='converted.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 8. PowerPoint to PDF
@app.route('/powerpoint-to-pdf', methods=['POST'])
def powerpoint_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PowerPoint file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        prs = Presentation(file)
        pdf = FPDF()
        
        for i, slide in enumerate(prs.slides):
            pdf.add_page()
            pdf.set_font("Arial", size=16)
            pdf.cell(0, 10, f"Slide {i + 1}", ln=True)
            pdf.ln(5)
            
            pdf.set_font("Arial", size=12)
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.encode('latin1', 'replace').decode('latin1')
                    pdf.multi_cell(0, 8, text)
                    pdf.ln(3)
        
        output = io.BytesIO()
        pdf_content = pdf.output(dest='S').encode('latin1')
        output.write(pdf_content)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='converted.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 9. Excel to PDF
@app.route('/excel-to-pdf', methods=['POST'])
def excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No Excel file uploaded'}), 400
    
    file = request.files['file']
    
    try:
        df = pd.read_excel(file)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=10)
        
        # Add headers
        for col in df.columns:
            pdf.cell(30, 10, str(col)[:10], 1)
        pdf.ln()
        
        # Add data rows
        for index, row in df.iterrows():
            for value in row:
                cell_value = str(value)[:10] if pd.notna(value) else ""
                pdf.cell(30, 10, cell_value, 1)
            pdf.ln()
            
            # Prevent page overflow
            if pdf.get_y() > 250:
                pdf.add_page()
        
        output = io.BytesIO()
        pdf_content = pdf.output(dest='S').encode('latin1')
        output.write(pdf_content)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='converted.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 10. Edit PDF (Add text/watermark)
@app.route('/edit-pdf', methods=['POST'])
def edit_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    edit_text = request.form.get('text', 'Sample Text')
    x_pos = int(request.form.get('x', 100))
    y_pos = int(request.form.get('y', 100))
    
    try:
        # Create overlay with text
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.drawString(x_pos, y_pos, edit_text)
        can.save()
        
        # Move to beginning of BytesIO buffer
        packet.seek(0)
        new_pdf = PdfReader(packet)
        
        # Read existing PDF
        existing_pdf = PdfReader(file)
        output = PdfWriter()
        
        # Add text to first page
        page = existing_pdf.pages[0]
        page.merge_page(new_pdf.pages[0])
        output.add_page(page)
        
        # Add remaining pages
        for i in range(1, len(existing_pdf.pages)):
            output.add_page(existing_pdf.pages[i])
        
        output_stream = io.BytesIO()
        output.write(output_stream)
        output_stream.seek(0)
        
        return send_file(output_stream, mimetype='application/pdf', download_name='edited.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 11. PDF to JPG
@app.route('/pdf-to-jpg', methods=['POST'])
def pdf_to_jpg():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    page_number = int(request.form.get('page', 1)) - 1
    
    try:
        temp_file = get_temp_filename('pdf')
        file.save(temp_file)
        
        doc = fitz.open(temp_file)
        
        if page_number >= len(doc):
            return jsonify({'error': 'Page number out of range'}), 400
        
        page = doc.load_page(page_number)
        mat = fitz.Matrix(2, 2)  # 2x zoom for better quality
        pix = page.get_pixmap(matrix=mat)
        
        img_data = pix.tobytes("jpeg")
        output = io.BytesIO(img_data)
        output.seek(0)
        
        doc.close()
        os.remove(temp_file)
        
        return send_file(output, mimetype='image/jpeg', download_name=f'page_{page_number+1}.jpg')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 12. JPG to PDF
@app.route('/jpg-to-pdf', methods=['POST'])
def jpg_to_pdf():
    if 'files' not in request.files:
        return jsonify({'error': 'No image files uploaded'}), 400
    
    files = request.files.getlist('files')
    
    try:
        pdf = FPDF()
        
        for file in files:
            if file and file.filename.lower().endswith(('.jpg', '.jpeg', '.png')):
                # Save image temporarily
                temp_img = get_temp_filename('jpg')
                file.save(temp_img)
                
                # Add to PDF
                pdf.add_page()
                pdf.image(temp_img, x=10, y=10, w=190)
                
                os.remove(temp_img)
        
        output = io.BytesIO()
        pdf_content = pdf.output(dest='S').encode('latin1')
        output.write(pdf_content)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='images.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 13. Sign PDF (Simple signature)
@app.route('/sign-pdf', methods=['POST'])
def sign_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    signature_text = request.form.get('signature', 'Digital Signature')
    
    try:
        # Create signature overlay
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFont("Helvetica-Bold", 12)
        can.drawString(400, 50, signature_text)
        can.drawString(400, 35, datetime.now().strftime("%Y-%m-%d %H:%M:%S"))
        can.save()
        
        packet.seek(0)
        signature_pdf = PdfReader(packet)
        
        existing_pdf = PdfReader(file)
        output = PdfWriter()
        
        # Add signature to last page
        for i, page in enumerate(existing_pdf.pages):
            if i == len(existing_pdf.pages) - 1:  # Last page
                page.merge_page(signature_pdf.pages[0])
            output.add_page(page)
        
        output_stream = io.BytesIO()
        output.write(output_stream)
        output_stream.seek(0)
        
        return send_file(output_stream, mimetype='application/pdf', download_name='signed.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 14. Watermark PDF
@app.route('/watermark-pdf', methods=['POST'])
def watermark_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    watermark_text = request.form.get('watermark', 'CONFIDENTIAL')
    
    try:
        # Create watermark
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.saveState()
        can.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
        can.setFont("Helvetica-Bold", 50)
        can.rotate(45)
        can.drawString(200, 0, watermark_text)
        can.restoreState()
        can.save()
        
        packet.seek(0)
        watermark_pdf = PdfReader(packet)
        
        existing_pdf = PdfReader(file)
        output = PdfWriter()
        
        # Apply watermark to all pages
        for page in existing_pdf.pages:
            page.merge_page(watermark_pdf.pages[0])
            output.add_page(page)
        
        output_stream = io.BytesIO()
        output.write(output_stream)
        output_stream.seek(0)
        
        return send_file(output_stream, mimetype='application/pdf', download_name='watermarked.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 15. Rotate PDF
@app.route('/rotate-pdf', methods=['POST'])
def rotate_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    rotation = int(request.form.get('rotation', 90))
    
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.rotate(rotation)
            writer.add_page(page)
        
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='rotated.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 16. HTML to PDF
@app.route('/html-to-pdf', methods=['POST'])
def html_to_pdf():
    html_content = request.form.get('html', '')
    
    if not html_content:
        return jsonify({'error': 'No HTML content provided'}), 400
    
    try:
        # Use pdfkit to convert HTML to PDF
        options = {
            'page-size': 'A4',
            'margin-top': '0.75in',
            'margin-right': '0.75in',
            'margin-bottom': '0.75in',
            'margin-left': '0.75in',
            'encoding': "UTF-8",
            'no-outline': None
        }
        
        pdf_content = pdfkit.from_string(html_content, False, options=options)
        output = io.BytesIO(pdf_content)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='converted.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 17. Unlock PDF (Remove password - basic implementation)
@app.route('/unlock-pdf', methods=['POST'])
def unlock_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    password = request.form.get('password', '')
    
    try:
        reader = PdfReader(file)
        
        if reader.is_encrypted:
            if not password:
                return jsonify({'error': 'Password required for encrypted PDF'}), 400
            
            if not reader.decrypt(password):
                return jsonify({'error': 'Invalid password'}), 400
        
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='unlocked.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 18. Protect PDF (Add password)
@app.route('/protect-pdf', methods=['POST'])
def protect_pdf():
    if 'file' not in request.files:
        return jsonify({'error': 'No PDF file uploaded'}), 400
    
    file = request.files['file']
    password = request.form.get('password', 'default123')
    
    try:
        reader = PdfReader(file)
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        writer.encrypt(password)
        
        output = io.BytesIO()
        writer.write(output)
        output.seek(0)
        
        return send_file(output, mimetype='application/pdf', download_name='protected.pdf')
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Health check endpoint
@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'message': 'Document conversion API is running'})

# List available conversions
@app.route('/conversions', methods=['GET'])
def list_conversions():
    conversions = {
        'PDF Operations': [
            'merge-pdf', 'split-pdf', 'compress-pdf', 'rotate-pdf', 
            'watermark-pdf', 'sign-pdf', 'edit-pdf', 'unlock-pdf', 'protect-pdf'
        ],
        'PDF Conversions': [
            'pdf-to-word', 'pdf-to-powerpoint', 'pdf-to-excel', 'pdf-to-jpg'
        ],
        'To PDF Conversions': [
            'word-to-pdf', 'powerpoint-to-pdf', 'excel-to-pdf', 
            'jpg-to-pdf', 'html-to-pdf'
        ],
        'Image Conversions': [
            'convert-image'
        ]
    }
    return jsonify(conversions)

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080, debug=True)